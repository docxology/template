"""Real-render tests for infrastructure.rendering.pptx_deck (no mocks).

Skips cleanly if the opt-in `rendering-pptx` dependency group
(`uv sync --group rendering-pptx`) is not installed, matching the pattern the
rest of the repo uses for optional-dependency modules (e.g. discopy tests).
"""

from __future__ import annotations

from pathlib import Path

import pdfplumber
import pytest

pytest.importorskip("pptx", reason="python-pptx opt-in group not installed (uv sync --group rendering-pptx)")

from pptx import Presentation  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE  # noqa: E402

from infrastructure.core.exceptions import RenderingError  # noqa: E402
from infrastructure.rendering.pptx_deck import render_pptx  # noqa: E402
from infrastructure.rendering.slide_deck import (  # noqa: E402
    CONTENT_PROTECTED_FLOOR_PT,
    DIAGRAM_FIGURE_BOTTOM_PT,
    DIAGRAM_FIGURE_MAX_HEIGHT_PT,
    DIAGRAM_FIGURE_MAX_WIDTH_PT,
    DIAGRAM_FIGURE_TOP_FROM_TOP_PT,
    PAGE_SIZE,
    SECTION_RULE_HEIGHT_PT,
    SECTION_RULE_TOP_FROM_TOP_PT,
    SECTION_TITLE_RULE_GAP_PT,
    DeckContent,
    DeckTheme,
    SLIDE_TEXT_WIDTH_PT,
    Slide,
    fit_helvetica_bold_single_line_font_size,
    plan_content_slide_layout,
    plan_diagram_figure_layout,
    render_pdf,
)


def _make_deck(n_slides: int) -> DeckContent:
    slides = tuple(
        Slide(title=f"Slide {i}", bullets=(f"Point {i}.a", f"Point {i}.b"), kind="content") for i in range(n_slides)
    )
    return DeckContent(title="Test Deck", subtitle="A subtitle", slides=slides)


def test_render_pptx_rejects_empty_deck(tmp_path: Path):
    empty_deck = DeckContent(title="Empty")
    with pytest.raises(RenderingError):
        render_pptx(empty_deck, tmp_path / "empty.pptx")


def test_render_pptx_produces_real_file_with_expected_slide_count(tmp_path: Path):
    deck = _make_deck(4)
    output = render_pptx(deck, tmp_path / "deck.pptx")

    assert output.is_file()
    assert output.stat().st_size > 1000

    prs = Presentation(str(output))
    assert len(prs.slides._sldIdLst) == 4 + 1  # synthesized title + 4 content slides


def test_render_pptx_is_byte_identical_for_identical_decks(tmp_path: Path):
    deck = _make_deck(2)

    first = render_pptx(deck, tmp_path / "first.pptx")
    second = render_pptx(deck, tmp_path / "second.pptx")

    assert first.read_bytes() == second.read_bytes()


def test_render_pptx_slide_count_matches_render_pdf_page_count(tmp_path: Path):
    from pypdf import PdfReader

    deck = _make_deck(6)
    pdf_path = render_pdf(deck, tmp_path / "deck.pdf")
    pptx_path = render_pptx(deck, tmp_path / "deck.pptx")

    pdf_page_count = len(PdfReader(str(pdf_path)).pages)
    prs = Presentation(str(pptx_path))
    pptx_slide_count = len(prs.slides._sldIdLst)

    assert pdf_page_count == pptx_slide_count


def test_render_pptx_text_contains_slide_titles_and_bullets(tmp_path: Path):
    deck = _make_deck(2)
    output = render_pptx(deck, tmp_path / "deck_text.pptx")

    prs = Presentation(str(output))
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = "\n".join(all_text)

    assert "Slide 0" in joined
    assert "Slide 1" in joined
    assert "Point 0.a" in joined


def test_render_pptx_no_synthesized_title_slide_when_first_slide_is_title(tmp_path: Path):
    slides = (
        Slide(title="Deck Title", kind="title"),
        Slide(title="Content", bullets=("a", "b")),
    )
    deck = DeckContent(title="Deck Title", slides=slides)
    output = render_pptx(deck, tmp_path / "deck_with_title_slide.pptx")

    prs = Presentation(str(output))
    assert len(prs.slides._sldIdLst) == 2


def test_render_pptx_notes_are_attached(tmp_path: Path):
    slides = (Slide(title="Deck Title", kind="title", notes="Speaker note here"),)
    deck = DeckContent(title="Deck Title", slides=slides)
    output = render_pptx(deck, tmp_path / "deck_notes.pptx")

    prs = Presentation(str(output))
    slide = next(iter(prs.slides))
    assert "Speaker note here" in slide.notes_slide.notes_text_frame.text


def test_render_pptx_figure_is_embedded(tmp_path: Path):
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig_path = tmp_path / "fig.png"
    fig, ax = plt.subplots()
    ax.plot([0, 1], [0, 1])
    fig.savefig(fig_path)
    plt.close(fig)

    slides = (Slide(title="With figure", bullets=("a",), figure_path=fig_path),)
    deck = DeckContent(title="Deck", slides=slides)
    output = render_pptx(deck, tmp_path / "deck_with_figure.pptx")

    prs = Presentation(str(output))
    slide = list(prs.slides)[-1]
    picture_shapes = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE == 13
    assert len(picture_shapes) == 1


def test_render_pptx_content_figure_is_below_estimated_bullet_box(tmp_path: Path):
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig_path = tmp_path / "flow.png"
    fig, ax = plt.subplots()
    ax.plot([0, 1], [0, 1])
    fig.savefig(fig_path)
    plt.close(fig)

    deck = DeckContent(
        title="Deck",
        slides=(
            Slide(
                title="Long content",
                bullets=(
                    "A deliberately long bullet that consumes multiple estimated lines and must reserve figure space.",
                    "A second long bullet keeps the regression focused on non-overlap rather than a one-line special case.",
                ),
                figure_path=fig_path,
            ),
        ),
    )
    output = render_pptx(deck, tmp_path / "flow.pptx")
    content_slide = list(Presentation(str(output)).slides)[-1]
    body = next(
        shape for shape in content_slide.shapes if shape.has_text_frame and "deliberately long bullet" in shape.text
    )
    picture = next(shape for shape in content_slide.shapes if shape.shape_type == 13)
    assert body.top + body.height <= picture.top


def test_render_pptx_stat_slide_shows_value_and_label(tmp_path: Path):
    deck = DeckContent(
        title="Deck",
        slides=(Slide(title="Proof", kind="stat", stat_value="89 tests", stat_label="91% coverage"),),
    )
    output = render_pptx(deck, tmp_path / "stat_deck.pptx")

    prs = Presentation(str(output))
    texts = [shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame]
    joined = "\n".join(texts)
    assert "89 tests" in joined
    assert "91% coverage" in joined


def test_render_pptx_quote_slide_shows_quote_and_attribution(tmp_path: Path):
    deck = DeckContent(
        title="Deck",
        slides=(Slide(title="", kind="quote", quote_text="A real quote.", quote_attribution="Some Source"),),
    )
    output = render_pptx(deck, tmp_path / "quote_deck.pptx")

    prs = Presentation(str(output))
    texts = [shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame]
    joined = "\n".join(texts)
    assert "A real quote." in joined
    assert "Some Source" in joined


def test_render_pptx_diagram_slide_embeds_figure(tmp_path: Path):
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig_path = tmp_path / "fig.png"
    fig, ax = plt.subplots()
    ax.plot([0, 1], [1, 0])
    fig.savefig(fig_path)
    plt.close(fig)

    deck = DeckContent(title="Deck", slides=(Slide(title="Architecture", kind="diagram", figure_path=fig_path),))
    output = render_pptx(deck, tmp_path / "diagram_deck.pptx")

    prs = Presentation(str(output))
    slide = list(prs.slides)[-1]
    picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
    assert len(picture_shapes) == 1


@pytest.mark.parametrize(
    ("image_size", "bound_dimension"),
    (((200, 1200), "height"), ((1200, 200), "width")),
)
def test_pdf_and_pptx_fit_diagram_to_same_protected_geometry(
    tmp_path: Path,
    image_size: tuple[int, int],
    bound_dimension: str,
):
    """Portrait and landscape images must fit by their limiting dimension.

    This is the real counterexample for the former PPTX width-only placement:
    at eight inches wide a 1:6 image was more than forty inches tall.
    """
    from PIL import Image

    figure_path = tmp_path / f"{bound_dimension}-bound-diagram.png"
    Image.new("RGB", image_size, "white").save(figure_path)
    slide = Slide(
        title="Tall evidence",
        kind="diagram",
        figure_path=figure_path,
        source="CLAUDE.md",
        qr_url="https://github.com/org/repo/blob/main/slide.md",
    )
    deck = DeckContent(title="Deck", slides=(slide,))
    expected = plan_diagram_figure_layout(figure_path)

    pdf_path = render_pdf(
        deck,
        tmp_path / "tall-diagram.pdf",
        source_base_url="https://github.com/org/repo/blob/main/",
    )
    pptx_path = render_pptx(
        deck,
        tmp_path / "tall-diagram.pptx",
        source_base_url="https://github.com/org/repo/blob/main/",
    )

    with pdfplumber.open(pdf_path) as document:
        diagram_page = document.pages[1]
        pdf_figure = next(image for image in diagram_page.images if image["srcsize"] == image_size)
        assert pdf_figure["x0"] == pytest.approx(expected.left_pt, abs=0.01)
        assert pdf_figure["y0"] == pytest.approx(expected.bottom_pt, abs=0.01)
        assert pdf_figure["width"] == pytest.approx(expected.width_pt, abs=0.01)
        assert pdf_figure["height"] == pytest.approx(expected.height_pt, abs=0.01)
        assert pdf_figure["top"] == pytest.approx(expected.top_pt, abs=0.01)

    prs = Presentation(str(pptx_path))
    diagram_slide = list(prs.slides)[1]
    pictures = [shape for shape in diagram_slide.shapes if shape.shape_type == 13]
    pptx_figure = max(pictures, key=lambda shape: shape.width * shape.height)
    assert pptx_figure.left / 12_700 == pytest.approx(expected.left_pt, abs=0.01)
    assert pptx_figure.top / 12_700 == pytest.approx(expected.top_pt, abs=0.01)
    assert pptx_figure.width / 12_700 == pytest.approx(expected.width_pt, abs=0.01)
    assert pptx_figure.height / 12_700 == pytest.approx(expected.height_pt, abs=0.01)
    pptx_figure_bottom_clearance = prs.slide_height - (pptx_figure.top + pptx_figure.height)
    assert pptx_figure_bottom_clearance >= DIAGRAM_FIGURE_BOTTOM_PT * 12_700 - 2
    assert expected.top_pt >= DIAGRAM_FIGURE_TOP_FROM_TOP_PT
    assert expected.width_pt <= DIAGRAM_FIGURE_MAX_WIDTH_PT
    assert expected.height_pt <= DIAGRAM_FIGURE_MAX_HEIGHT_PT
    assert expected.top_pt + expected.height_pt <= PAGE_SIZE[1] - DIAGRAM_FIGURE_BOTTOM_PT
    assert expected.width_pt / expected.height_pt == pytest.approx(image_size[0] / image_size[1])
    if bound_dimension == "height":
        assert expected.height_pt == pytest.approx(DIAGRAM_FIGURE_MAX_HEIGHT_PT)
    else:
        assert expected.width_pt == pytest.approx(DIAGRAM_FIGURE_MAX_WIDTH_PT)


def test_pdf_and_pptx_section_title_rule_geometry_does_not_intersect(tmp_path: Path):
    title = "Scientific integrity, by construction"
    deck = DeckContent(title="Deck", slides=(Slide(title=title, kind="section"),))
    pdf_path = render_pdf(deck, tmp_path / "section.pdf")
    pptx_path = render_pptx(deck, tmp_path / "section.pptx")

    with pdfplumber.open(pdf_path) as document:
        page = document.pages[1]
        pdf_rule = next(
            rect
            for rect in page.rects
            if rect["height"] == pytest.approx(SECTION_RULE_HEIGHT_PT) and rect["width"] == pytest.approx(640.8)
        )
        expected_title_words = title.split()
        title_words = [word for word in page.extract_words() if word["text"] in expected_title_words]
        assert pdf_rule["top"] == pytest.approx(SECTION_RULE_TOP_FROM_TOP_PT)
        assert [word["text"] for word in title_words] == expected_title_words
        assert max(word["bottom"] for word in title_words) < pdf_rule["top"]

    prs = Presentation(str(pptx_path))
    section_slide = list(prs.slides)[1]
    title_frame, _run = _pptx_title_frame_and_run(section_slide, title)
    title_box = title_frame._parent
    rule = next(
        shape for shape in section_slide.shapes if shape.height == pytest.approx(SECTION_RULE_HEIGHT_PT * 12_700, abs=1)
    )
    assert rule.top == pytest.approx(SECTION_RULE_TOP_FROM_TOP_PT * 12_700, abs=1)
    assert title_box.top + title_box.height <= rule.top
    assert rule.top - (title_box.top + title_box.height) == pytest.approx(
        SECTION_TITLE_RULE_GAP_PT * 12_700,
        abs=2,
    )
    assert title_frame.word_wrap is False
    assert title_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert title_frame.vertical_anchor == MSO_ANCHOR.TOP
    assert title_frame.margin_left == title_frame.margin_right == 0
    assert title_frame.margin_top == title_frame.margin_bottom == 0


def test_render_pptx_source_citation_becomes_clickable_link(tmp_path: Path):
    deck = DeckContent(title="Deck", slides=(Slide(title="Fact slide", bullets=("x",), source="CLAUDE.md"),))
    output = render_pptx(deck, tmp_path / "cited_deck.pptx", source_base_url="https://github.com/org/repo/blob/main/")

    prs = Presentation(str(output))
    content_slide = list(prs.slides)[-1]  # slide 0 is the synthesized title slide
    found_hyperlink = None
    found_text = False
    for shape in content_slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "Source: CLAUDE.md" in run.text:
                    found_text = True
                if run.hyperlink.address:
                    found_hyperlink = run.hyperlink.address

    assert found_text
    assert found_hyperlink == "https://github.com/org/repo/blob/main/CLAUDE.md"


def test_render_pptx_custom_theme_changes_output(tmp_path: Path):
    """Different themes must actually change the rendered artifact (not a no-op parameter)."""
    deck = DeckContent(title="Deck", slides=(Slide(title="Title", kind="title"),))
    default_output = render_pptx(deck, tmp_path / "default_theme.pptx")
    custom_theme = DeckTheme(
        black="#123456", white="#FFFFFF", highlight_1="#ABCDEF", highlight_2="#ABCDEF", highlight_3="#ABCDEF"
    )
    custom_output = render_pptx(deck, tmp_path / "custom_theme.pptx", theme=custom_theme)
    assert default_output.read_bytes() != custom_output.read_bytes()


def test_render_pptx_empty_title_does_not_crash(tmp_path: Path):
    """Regression test: python-pptx leaves a paragraph with zero runs when
    text is set to "", so a naive `.text = title` + `runs[0]` access used to
    raise IndexError on any empty-string title (title/section/content/diagram
    slide kinds all accept one)."""
    deck = DeckContent(title="", slides=(Slide(title="", kind="title"), Slide(title="", bullets=("x",))))
    output = render_pptx(deck, tmp_path / "empty_title.pptx")
    assert output.is_file()


def test_render_pptx_missing_figure_warns_and_renders_without_it(tmp_path: Path, caplog):
    missing_path = tmp_path / "does_not_exist.png"
    deck = DeckContent(title="Deck", slides=(Slide(title="Content", bullets=("x",), figure_path=missing_path),))
    with caplog.at_level("WARNING"):
        output = render_pptx(deck, tmp_path / "missing_figure.pptx")
    assert output.is_file()
    assert any("does not exist" in record.message for record in caplog.records)


def test_render_pptx_missing_figure_on_diagram_slide_warns_and_renders(tmp_path: Path, caplog):
    missing_path = tmp_path / "does_not_exist.png"
    deck = DeckContent(title="Deck", slides=(Slide(title="Diagram", kind="diagram", figure_path=missing_path),))
    with caplog.at_level("WARNING"):
        output = render_pptx(deck, tmp_path / "missing_diagram_figure.pptx")
    assert output.is_file()
    assert any("does not exist" in record.message for record in caplog.records)


def test_render_pptx_draws_clickable_qr_code_when_qr_url_set(tmp_path: Path):
    qr_target = "https://github.com/org/repo/blob/main/slide.md"
    deck = DeckContent(title="Deck", slides=(Slide(title="Fact slide", bullets=("x",), qr_url=qr_target),))
    output = render_pptx(deck, tmp_path / "qr_deck.pptx")

    prs = Presentation(str(output))
    content_slide = list(prs.slides)[-1]
    picture_shapes = [s for s in content_slide.shapes if s.shape_type == 13]
    assert len(picture_shapes) == 1
    assert picture_shapes[0].click_action.hyperlink.address == qr_target


def test_render_pptx_no_qr_picture_when_qr_url_absent(tmp_path: Path):
    deck = DeckContent(title="Deck", slides=(Slide(title="No QR", bullets=("x",)),))
    output = render_pptx(deck, tmp_path / "no_qr_deck.pptx")
    prs = Presentation(str(output))
    content_slide = list(prs.slides)[-1]
    picture_shapes = [s for s in content_slide.shapes if s.shape_type == 13]
    assert len(picture_shapes) == 0


def test_render_pptx_qr_and_figure_coexist_as_two_pictures(tmp_path: Path):
    """A diagram slide has its own figure picture; the QR must be a second,
    distinct picture shape, not clobber the figure."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig_path = tmp_path / "fig.png"
    fig, ax = plt.subplots()
    ax.plot([0, 1], [1, 0])
    fig.savefig(fig_path)
    plt.close(fig)

    deck = DeckContent(
        title="Deck",
        slides=(
            Slide(
                title="Diagram",
                kind="diagram",
                figure_path=fig_path,
                qr_url="https://github.com/org/repo/blob/main/slide.md",
            ),
        ),
    )
    output = render_pptx(deck, tmp_path / "diagram_with_qr.pptx")
    prs = Presentation(str(output))
    content_slide = list(prs.slides)[-1]
    picture_shapes = [s for s in content_slide.shapes if s.shape_type == 13]
    assert len(picture_shapes) == 2


def _pdf_text_font_sizes(page, target: str) -> list[float]:
    sizes: list[float] = []

    def visit(text, _cm, _tm, _font, font_size):
        if text.strip() == target:
            sizes.append(float(font_size))

    page.extract_text(visitor_text=visit)
    return sizes


def _pptx_title_frame_and_run(slide, target: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text == target:
                    return shape.text_frame, run
    raise AssertionError(f"missing title run: {target}")


def test_pdf_and_pptx_match_all_eight_regressed_title_sizes(tmp_path: Path):
    from pypdf import PdfReader

    cases = (
        ("The full roster, not one cherry-picked example", "content", 24),
        ("Scientific integrity, by construction", "section", 33),
        ("What's actually inside infrastructure/", "section", 33),
        ("Why this is a science-integrity problem, not just a tooling problem", "content", 20),
        ("The full roster, not one cherry-picked example", "content", 24),
        ("Scientific integrity, by construction", "section", 33),
        ("What's actually inside infrastructure/", "section", 33),
        ("Where this could go next for your organization", "content", 24),
    )
    deck = DeckContent(
        title="Deck",
        slides=tuple(Slide(title=title, kind=kind, bullets=("Short body.",)) for title, kind, _ in cases),
    )
    pdf = render_pdf(deck, tmp_path / "titles.pdf")
    pptx = render_pptx(deck, tmp_path / "titles.pptx")
    pdf_pages = PdfReader(str(pdf)).pages
    pptx_slides = list(Presentation(str(pptx)).slides)

    for index, (title, kind, expected_size) in enumerate(cases, start=1):
        assert _pdf_text_font_sizes(pdf_pages[index], title) == [float(expected_size)]
        pptx_frame, pptx_run = _pptx_title_frame_and_run(pptx_slides[index], title)
        assert pptx_run.font.size.pt == expected_size
        assert pptx_run.font.name == "Helvetica"
        assert pptx_frame.word_wrap is False
        assert pptx_frame.auto_size == MSO_AUTO_SIZE.NONE
        if kind == "section":
            assert pptx_frame.margin_left == 0
            assert pptx_frame.margin_right == 0
        else:
            assert pptx_frame.margin_left == 502_920  # 0.55 inch
            assert pptx_frame.margin_right == 502_920


def test_pptx_section_title_uses_exact_shared_text_width_without_wrapping(tmp_path: Path):
    title = "W" * 21
    assert (
        fit_helvetica_bold_single_line_font_size(
            title,
            max_width_pt=SLIDE_TEXT_WIDTH_PT,
            start_size_pt=33,
        )
        == 32
    )
    deck = DeckContent(title="Deck", slides=(Slide(title=title, kind="section"),))
    output = render_pptx(deck, tmp_path / "section-title-width.pptx")
    slide = list(Presentation(str(output)).slides)[-1]
    frame, run = _pptx_title_frame_and_run(slide, title)

    assert run.font.size.pt == 32
    assert frame.word_wrap is False
    assert frame.auto_size == MSO_AUTO_SIZE.NONE
    assert frame.margin_left == 0
    assert frame.margin_right == 0
    assert frame._parent.width == pytest.approx(SLIDE_TEXT_WIDTH_PT * 12_700, abs=1)


def test_render_pptx_body_overflow_preserves_existing_target(tmp_path: Path):
    output = tmp_path / "sentinel.pptx"
    output.write_bytes(b"sentinel-body")
    deck = DeckContent(title="Deck", slides=(Slide(title="Overflow", bullets=(("many words " * 200).strip(),)),))
    with pytest.raises(RenderingError, match="protected footer/QR band"):
        render_pptx(deck, output)
    assert output.read_bytes() == b"sentinel-body"


def test_render_pptx_title_overflow_preserves_existing_target(tmp_path: Path):
    output = tmp_path / "sentinel.pptx"
    output.write_bytes(b"sentinel-title")
    deck = DeckContent(title="Deck", slides=(Slide(title="W" * 500, kind="diagram"),))
    with pytest.raises(RenderingError, match="minimum font size"):
        render_pptx(deck, output)
    assert output.read_bytes() == b"sentinel-title"


def test_pptx_body_consumes_shared_boundary_plan_above_protected_floor(tmp_path: Path):
    planned_slide = Slide(
        title="Planned body",
        bullets=("word", ("word " * 141).strip()),
    )
    layout = plan_content_slide_layout(planned_slide)
    assert sum(len(lines) for lines in layout.bullet_lines) == 12
    assert layout.last_glyph_bottom_pt == pytest.approx(64.794)
    deck = DeckContent(title="Deck", slides=(planned_slide,))
    output = render_pptx(deck, tmp_path / "planned.pptx")
    prs = Presentation(str(output))
    slide = list(prs.slides)[-1]
    body = next(shape for shape in slide.shapes if shape.has_text_frame and "word" in shape.text)
    assert "\v" in body.text
    assert body.text_frame.word_wrap is False
    assert body.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert body.text_frame.vertical_anchor == MSO_ANCHOR.TOP
    assert body.top == pytest.approx(layout.body_top_pt * 12_700, abs=1)
    assert prs.slide_height - (body.top + body.height) >= CONTENT_PROTECTED_FLOOR_PT * 12700
    assert prs.slide_height - (body.top + body.height) == pytest.approx(
        layout.last_glyph_bottom_pt * 12_700,
        abs=2,
    )
    assert body.text_frame.paragraphs[0].line_spacing.pt == pytest.approx(21.0)
    assert body.text_frame.paragraphs[0].space_after.pt == pytest.approx(8.28)


def test_render_pptx_content_slide_with_long_title_does_not_crash(tmp_path: Path):
    long_title = "Why this is a science-integrity problem, not just a tooling problem"
    deck = DeckContent(title="Deck", slides=(Slide(title=long_title, bullets=("x",)),))
    output = render_pptx(deck, tmp_path / "long_title.pptx")
    assert output.is_file()
