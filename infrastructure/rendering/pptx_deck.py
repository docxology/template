"""PPTX renderer sharing :class:`infrastructure.rendering.slide_deck.DeckContent`.

The first PowerPoint-generation capability in this repository (confirmed
absent, repo-wide, before this module). Deliberately mirrors
:func:`infrastructure.rendering.slide_deck.render_pdf` slide-for-slide — same
:class:`~infrastructure.rendering.slide_deck.DeckTheme`, same font-size
constants (imported directly from `slide_deck`, not re-declared, so the two
renderers cannot drift apart the way they once did), same slide `kind`
dispatch (title/section/content/stat/quote/diagram), same `source_base_url`
deep-link behavior and per-slide QR support — so a PDF and a PPTX built from
the same :class:`DeckContent` have identical content slide counts and
identical bullet text.

``python-pptx`` is an opt-in dependency (``uv sync --group rendering-pptx``),
not part of the default dependency groups — importing this module without
that group installed raises :class:`ImportError` with a clear message rather
than a bare traceback.
"""

from __future__ import annotations

import io
import os
from pathlib import Path
import tempfile
from typing import Any
from zipfile import ZipFile, ZipInfo

from infrastructure.core.exceptions import RenderingError
from infrastructure.core.logging.utils import get_logger
from infrastructure.rendering.slide_deck import (
    CONTENT_BODY_FONT_SIZE,
    CONTENT_BODY_LINE_HEIGHT_PT,
    CONTENT_HEADER_FONT_SIZE,
    DEFAULT_THEME,
    DIAGRAM_FIGURE_BOTTOM_PT,
    DIAGRAM_HEADER_FONT_SIZE,
    QUOTE_ATTRIBUTION_FONT_SIZE,
    QUOTE_FONT_SIZE,
    SECTION_FONT_SIZE,
    SECTION_RULE_HEIGHT_PT,
    SECTION_RULE_TOP_FROM_TOP_PT,
    SECTION_TITLE_BOX_HEIGHT_PT,
    SECTION_TITLE_BOX_TOP_PT,
    SOURCE_FOOTER_FONT_SIZE,
    STAT_LABEL_FONT_SIZE,
    STAT_TITLE_FONT_SIZE,
    STAT_VALUE_FONT_SIZE,
    SUBTITLE_FONT_SIZE,
    TITLE_FONT_SIZE,
    ContentSlideLayout,
    DeckContent,
    DeckTheme,
    Slide,
    SLIDE_TEXT_WIDTH_PT,
    fit_helvetica_bold_single_line_font_size,
    plan_diagram_figure_layout,
    source_url,
    validate_deck_layout,
)

logger = get_logger(__name__)

try:
    from pptx import Presentation
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    _PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover - exercised only when the opt-in group is absent
    _PPTX_AVAILABLE = False

_MSO_SHAPE_RECTANGLE = 1
_MSO_SHAPE_TYPE_PICTURE = 13
_QR_SIZE_INCHES = 0.62

_DETERMINISTIC_ZIP_TIMESTAMP = (1980, 1, 1, 0, 0, 0)


def is_pptx_available() -> bool:
    """Whether the opt-in ``python-pptx`` dependency group is installed.

    Importing this module never raises even when the group is absent (the
    ``pptx`` import above is caught internally) — callers that need to know
    ahead of a `render_pptx` call whether it will actually work (rather than
    catching the resulting `RenderingError`) should check this first.
    """
    return _PPTX_AVAILABLE


def render_pptx(
    deck: DeckContent, output_path: Path, *, theme: DeckTheme = DEFAULT_THEME, source_base_url: str = ""
) -> Path:
    """Render ``deck`` to a 16:9 ``.pptx`` file at ``output_path``.

    One slide per :class:`~infrastructure.rendering.slide_deck.Slide`, in the
    same order used by ``render_pdf`` — see that function's docstring for the
    shared front-matter convention and the diligence deep-link contract.
    Every slide whose `Slide.qr_url` is set gets a scannable, click-through
    QR code bottom-right, matching the PDF renderer.
    """
    if not _PPTX_AVAILABLE:
        raise RenderingError(
            "python-pptx is not installed. Install the opt-in dependency group "
            "with `uv sync --group rendering-pptx` before rendering PPTX output.",
            context={"deck_title": deck.title},
        )
    if not deck.slides:
        raise RenderingError("Cannot render a deck with zero slides", context={"deck_title": deck.title})
    layouts = validate_deck_layout(deck)

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    blank_layout = prs.slide_layouts[6]  # "Blank" layout present in the default template

    needs_title_slide = not deck.slides or deck.slides[0].kind != "title"
    if needs_title_slide:
        _add_title_slide(prs, blank_layout, deck.title, deck.subtitle, theme)

    for slide, content_layout in zip(deck.slides, layouts, strict=True):
        if slide.kind == "title":
            pptx_slide = _add_title_slide(prs, blank_layout, slide.title, deck.subtitle, theme, notes=slide.notes)
        elif slide.kind == "section":
            pptx_slide = _add_section_slide(prs, blank_layout, slide.title, theme, notes=slide.notes)
        elif slide.kind == "stat":
            pptx_slide = _add_stat_slide(prs, blank_layout, slide, theme)
        elif slide.kind == "quote":
            pptx_slide = _add_quote_slide(prs, blank_layout, slide, theme)
        elif slide.kind == "diagram":
            pptx_slide = _add_diagram_slide(prs, blank_layout, slide, theme)
        else:
            if content_layout is None:  # pragma: no cover - guarded by validate_deck_layout
                raise RenderingError("Missing content-slide layout", context={"slide_title": slide.title})
            pptx_slide = _add_content_slide(prs, blank_layout, slide, theme, content_layout)
        _add_source_footer(prs, pptx_slide, slide, theme, source_base_url)
        _add_qr_code(prs, pptx_slide, slide)

    prs.save(str(output_path))
    _normalize_pptx_archive(output_path)
    return output_path


def _normalize_pptx_archive(path: Path) -> None:
    """Rewrite an OOXML archive with stable ZIP metadata and member order.

    ``python-pptx`` emits stable XML and media bytes, but :mod:`zipfile` stamps
    every package member with the wall clock. That makes identical decks hash
    differently across runs even though their extracted contents match. PPTX
    uses ZIP only as a container, so fixed member timestamps preserve document
    semantics while making the published artifact byte-reproducible.
    """
    with ZipFile(path, "r") as source:
        members = [(info, source.read(info.filename)) for info in source.infolist()]

    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix=f".{path.name}.", suffix=".tmp", dir=path.parent, delete=False
        ) as handle:
            temporary_path = Path(handle.name)
        with ZipFile(temporary_path, "w") as target:
            for source_info, payload in members:
                target_info = ZipInfo(source_info.filename, _DETERMINISTIC_ZIP_TIMESTAMP)
                target_info.compress_type = source_info.compress_type
                target_info.comment = source_info.comment
                target_info.create_system = source_info.create_system
                target_info.external_attr = source_info.external_attr
                target_info.internal_attr = source_info.internal_attr
                target.writestr(target_info, payload)
        os.replace(temporary_path, path)
        temporary_path = None
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)


def _set_notes(pptx_slide: Any, notes: str) -> None:
    if notes:
        pptx_slide.notes_slide.notes_text_frame.text = notes


def _set_run_text(text_frame: Any, text: str) -> Any:
    """Set the first paragraph's text via `add_run()`, returning that run.

    ``text_frame.text = ""`` (and ``paragraph.text = ""``) both leave the
    paragraph with **zero runs** in python-pptx — a subsequent
    ``paragraphs[0].runs[0]`` then raises ``IndexError``. Every slide kind
    here may legitimately receive an empty string (e.g. a `quote`-kind
    slide's unused `title`), so every text-setting call site uses this
    helper instead of the fragile ``.text =`` + ``runs[0]`` idiom.
    """
    run = text_frame.paragraphs[0].add_run()
    run.text = text
    return run


def _add_source_footer(prs: Any, pptx_slide: Any, slide: Slide, theme: DeckTheme, source_base_url: str) -> None:
    if not slide.source:
        return
    url = source_url(slide.source, source_base_url)
    box = pptx_slide.shapes.add_textbox(Inches(0.55), prs.slide_height - Inches(0.35), Inches(6.0), Inches(0.3))
    run = _set_run_text(box.text_frame, f"Source: {slide.source}")
    run.font.size = Pt(SOURCE_FOOTER_FONT_SIZE)
    run.font.italic = True
    run.font.color.rgb = _rgb("8A8F99")
    if url:
        run.hyperlink.address = url


def _add_qr_code(prs: Any, pptx_slide: Any, slide: Slide) -> None:
    """Add a scannable + click-through QR code bottom-right, if `slide.qr_url` is set.

    Reuses `infrastructure.steganography.barcode_generators.generate_qr_code`
    (the same QR generator the steganography layer uses) rather than
    duplicating QR-encoding logic — matches `slide_deck.py`'s PDF placement.
    """
    if not slide.qr_url:
        return
    from infrastructure.steganography.barcode_generators import generate_qr_code

    png_bytes = generate_qr_code(slide.qr_url, box_size=3, border=1)
    size = Inches(_QR_SIZE_INCHES)
    x = prs.slide_width - Inches(0.55) - size
    y = prs.slide_height - Inches(0.74)
    picture = pptx_slide.shapes.add_picture(io.BytesIO(png_bytes), x, y, width=size, height=size)
    picture.click_action.hyperlink.address = slide.qr_url


def _add_title_slide(prs: Any, layout: Any, title: str, subtitle: str, theme: DeckTheme, notes: str = "") -> Any:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, prs, theme.black)

    accent = slide.shapes.add_shape(_MSO_SHAPE_RECTANGLE, 0, Inches(1.7), Inches(0.12), Inches(1.9))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme.highlight_1)
    accent.line.fill.background()
    accent.shadow.inherit = False

    # Title box is tall enough for a two-line wrap at TITLE_FONT_SIZE (a long
    # title is the common case, not the exception) so the fixed-position
    # subtitle box below it never overlaps — python-pptx text frames do not
    # auto-measure text extent, so this height must be generous by
    # construction rather than computed from the actual wrapped line count.
    box = slide.shapes.add_textbox(Inches(0.55), Inches(1.55), prs.slide_width - Inches(1.1), Inches(1.85))
    tf = box.text_frame
    tf.word_wrap = True
    run = _set_run_text(tf, title)
    run.font.size = Pt(TITLE_FONT_SIZE)
    run.font.bold = True
    run.font.color.rgb = _rgb(theme.white)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(3.6), prs.slide_width - Inches(1.1), Inches(0.9))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_run = _set_run_text(sub_tf, subtitle)
        sub_run.font.size = Pt(SUBTITLE_FONT_SIZE)
        sub_run.font.color.rgb = _rgb("C9C9C9")

    _set_notes(slide, notes)
    return slide


def _add_section_slide(prs: Any, layout: Any, title: str, theme: DeckTheme, notes: str = "") -> Any:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, prs, theme.white)

    rule = slide.shapes.add_shape(
        _MSO_SHAPE_RECTANGLE,
        Inches(0.55),
        Pt(SECTION_RULE_TOP_FROM_TOP_PT),
        prs.slide_width - Inches(1.1),
        Pt(SECTION_RULE_HEIGHT_PT),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(theme.highlight_1)
    rule.line.fill.background()
    rule.shadow.inherit = False

    box = slide.shapes.add_textbox(
        Inches(0.55),
        Pt(SECTION_TITLE_BOX_TOP_PT),
        prs.slide_width - Inches(1.1),
        Pt(SECTION_TITLE_BOX_HEIGHT_PT),
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    run = _set_run_text(tf, title)
    title_font_size = fit_helvetica_bold_single_line_font_size(
        title, max_width_pt=SLIDE_TEXT_WIDTH_PT, start_size_pt=SECTION_FONT_SIZE
    )
    run.font.size = Pt(title_font_size)
    run.font.name = "Helvetica"
    run.font.bold = True
    run.font.color.rgb = _rgb(theme.black)
    tf.paragraphs[0].line_spacing = Pt(title_font_size)

    _set_notes(slide, notes)
    return slide


def _add_content_slide(
    prs: Any,
    layout: Any,
    slide_content: Slide,
    theme: DeckTheme,
    content_layout: ContentSlideLayout,
) -> Any:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, prs, theme.white)

    header = slide.shapes.add_shape(_MSO_SHAPE_RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    header.fill.solid()
    header.fill.fore_color.rgb = _rgb(theme.black)
    header.line.fill.background()
    header.shadow.inherit = False
    header_tf = header.text_frame
    header_tf.word_wrap = False
    header_tf.auto_size = MSO_AUTO_SIZE.NONE
    header_run = _set_run_text(header_tf, slide_content.title)
    header_run.font.size = Pt(
        fit_helvetica_bold_single_line_font_size(
            slide_content.title,
            max_width_pt=SLIDE_TEXT_WIDTH_PT,
            start_size_pt=CONTENT_HEADER_FONT_SIZE,
        )
    )
    header_run.font.name = "Helvetica"
    header_run.font.bold = True
    header_run.font.color.rgb = _rgb(theme.white)
    header_tf.margin_left = Inches(0.55)
    header_tf.margin_right = Inches(0.55)

    body_top = Pt(content_layout.body_top_pt)
    body_height = Pt(content_layout.body_height_pt)
    if content_layout.figure_y_pt is not None and content_layout.figure_height_pt is not None:
        figure_top = 5.625 - (content_layout.figure_y_pt + content_layout.figure_height_pt) / 72
    body_box = slide.shapes.add_textbox(Inches(0.55), body_top, prs.slide_width - Inches(1.1), body_height)
    tf = body_box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (lines, baselines) in enumerate(
        zip(content_layout.bullet_lines, content_layout.line_baselines_pt, strict=True)
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "\v".join(lines)
        p.font.name = "Helvetica"
        p.font.size = Pt(CONTENT_BODY_FONT_SIZE)
        p.font.color.rgb = _rgb(theme.black)
        line_spacing = baselines[0] - baselines[1] if len(baselines) > 1 else CONTENT_BODY_LINE_HEIGHT_PT
        p.line_spacing = Pt(line_spacing)
        if i < len(content_layout.bullet_lines) - 1:
            next_baselines = content_layout.line_baselines_pt[i + 1]
            paragraph_gap = baselines[-1] - next_baselines[0] - line_spacing
            p.space_after = Pt(paragraph_gap)
        else:
            p.space_after = Pt(0)

    if slide_content.figure_path is not None:
        if slide_content.figure_path.is_file():
            if content_layout.figure_y_pt is None or content_layout.figure_height_pt is None:
                raise RenderingError("Missing content-figure layout", context={"slide_title": slide_content.title})
            slide.shapes.add_picture(
                str(slide_content.figure_path),
                Inches(1.5),
                Inches(figure_top),
                width=prs.slide_width - Inches(3.0),
                height=Pt(content_layout.figure_height_pt),
            )
        else:
            logger.warning(
                "Slide %r declares figure_path=%s but the file does not exist — "
                "rendering the slide without it rather than failing silently.",
                slide_content.title,
                slide_content.figure_path,
            )

    _set_notes(slide, slide_content.notes)
    return slide


def _add_stat_slide(prs: Any, layout: Any, slide_content: Slide, theme: DeckTheme) -> Any:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, prs, theme.white)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.4), prs.slide_width - Inches(1.1), Inches(0.6))
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = slide_content.title
    title_run.font.size = Pt(STAT_TITLE_FONT_SIZE)
    title_run.font.bold = True
    title_run.font.color.rgb = _rgb(theme.black)

    value = slide_content.stat_value or (slide_content.bullets[0] if slide_content.bullets else "")
    value_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.8), prs.slide_width - Inches(1.1), Inches(1.4))
    value_run = value_box.text_frame.paragraphs[0].add_run()
    value_run.text = value
    value_run.font.size = Pt(STAT_VALUE_FONT_SIZE)
    value_run.font.bold = True
    value_run.font.color.rgb = _rgb(theme.highlight_2)

    if slide_content.stat_label:
        label_box = slide.shapes.add_textbox(Inches(0.55), Inches(3.35), prs.slide_width - Inches(1.1), Inches(0.8))
        label_tf = label_box.text_frame
        label_tf.word_wrap = True
        label_run = label_tf.paragraphs[0].add_run()
        label_run.text = slide_content.stat_label
        label_run.font.size = Pt(STAT_LABEL_FONT_SIZE)
        label_run.font.color.rgb = _rgb(theme.black)

    return slide


def _add_quote_slide(prs: Any, layout: Any, slide_content: Slide, theme: DeckTheme) -> Any:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, prs, theme.black)

    accent = slide.shapes.add_shape(_MSO_SHAPE_RECTANGLE, Inches(0.55), Inches(1.7), Inches(0.55), Pt(6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme.highlight_3)
    accent.line.fill.background()
    accent.shadow.inherit = False

    quote_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.95), prs.slide_width - Inches(1.1), Inches(1.8))
    quote_tf = quote_box.text_frame
    quote_tf.word_wrap = True
    quote_run = quote_tf.paragraphs[0].add_run()
    quote_run.text = f"“{slide_content.quote_text}”"
    quote_run.font.size = Pt(QUOTE_FONT_SIZE)
    quote_run.font.italic = True
    quote_run.font.color.rgb = _rgb(theme.white)

    if slide_content.quote_attribution:
        attr_box = slide.shapes.add_textbox(Inches(0.55), Inches(3.8), prs.slide_width - Inches(1.1), Inches(0.5))
        attr_run = attr_box.text_frame.paragraphs[0].add_run()
        attr_run.text = f"— {slide_content.quote_attribution}"
        attr_run.font.size = Pt(QUOTE_ATTRIBUTION_FONT_SIZE)
        attr_run.font.bold = True
        attr_run.font.color.rgb = _rgb(theme.highlight_3)

    return slide


def _add_diagram_slide(prs: Any, layout: Any, slide_content: Slide, theme: DeckTheme) -> Any:
    slide = prs.slides.add_slide(layout)
    _fill_background(slide, prs, theme.white)

    header = slide.shapes.add_shape(_MSO_SHAPE_RECTANGLE, 0, 0, prs.slide_width, Inches(0.75))
    header.fill.solid()
    header.fill.fore_color.rgb = _rgb(theme.black)
    header.line.fill.background()
    header.shadow.inherit = False
    header_tf = header.text_frame
    header_tf.word_wrap = False
    header_tf.auto_size = MSO_AUTO_SIZE.NONE
    header_run = _set_run_text(header_tf, slide_content.title)
    header_run.font.size = Pt(
        fit_helvetica_bold_single_line_font_size(
            slide_content.title,
            max_width_pt=SLIDE_TEXT_WIDTH_PT,
            start_size_pt=DIAGRAM_HEADER_FONT_SIZE,
        )
    )
    header_run.font.name = "Helvetica"
    header_run.font.bold = True
    header_run.font.color.rgb = _rgb(theme.white)
    header_tf.margin_left = Inches(0.55)
    header_tf.margin_right = Inches(0.55)

    if slide_content.figure_path is not None:
        if slide_content.figure_path.is_file():
            figure_layout = plan_diagram_figure_layout(slide_content.figure_path)
            picture = slide.shapes.add_picture(
                str(slide_content.figure_path),
                Pt(figure_layout.left_pt),
                Pt(figure_layout.top_pt),
                width=Pt(figure_layout.width_pt),
                height=Pt(figure_layout.height_pt),
            )
            if prs.slide_height - (picture.top + picture.height) < Pt(DIAGRAM_FIGURE_BOTTOM_PT):
                raise RenderingError(
                    "Diagram figure enters the protected footer/QR band",
                    context={"slide_title": slide_content.title},
                )
        else:
            logger.warning(
                "Slide %r declares figure_path=%s but the file does not exist — "
                "rendering the slide without it rather than failing silently.",
                slide_content.title,
                slide_content.figure_path,
            )

    return slide


def _fill_background(slide: Any, prs: Any, hex_color: str) -> None:
    bg_shape = slide.shapes.add_shape(_MSO_SHAPE_RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = _rgb(hex_color)
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    # Send background rectangle behind subsequently-added shapes.
    spTree = slide.shapes._spTree
    spTree.remove(bg_shape._element)
    spTree.insert(2, bg_shape._element)


def _rgb(hex_color: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_color.lstrip("#"))
